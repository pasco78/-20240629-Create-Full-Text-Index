import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd

def extract_text_from_pdf(file_path):
    document = fitz.open(file_path)
    text = ""
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        text += page.get_text("text")
    return text

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()
