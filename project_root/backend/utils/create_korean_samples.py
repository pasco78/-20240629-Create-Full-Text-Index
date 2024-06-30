import sys
import os

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..')))

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
import openpyxl
from openpyxl import Workbook
from backend.utils.insert_data import insert_into_database
from backend.utils.extract_text import extract_text_from_pdf, extract_text_from_ppt
import os

def create_sample_pdf(file_path):
    c = canvas.Canvas(file_path, pagesize=letter)
    c.drawString(100, 750, "이것은 샘플 PDF 문서입니다.")
    c.drawString(100, 735, "테스트 목적으로 일부 샘플 텍스트가 포함되어 있습니다.")
    c.save()

def create_sample_ppt(file_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "샘플 프레젠테이션"
    subtitle.text = "테스트 목적으로 샘플 PPT 슬라이드입니다."
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title, content = slide.shapes.title, slide.placeholders[1]
    title.text = "슬라이드 2 제목"
    content.text = "이 슬라이드는 더 많은 샘플 텍스트를 포함하고 있습니다."
    prs.save(file_path)

def create_sample_excel(file_path):
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "샘플 엑셀 파일"
    ws['A2'] = "테스트 목적으로 일부 샘플 텍스트가 포함되어 있습니다."
    wb.save(file_path)

def create_sample_txt(file_path):
    with open(file_path, "w", encoding="utf-8") as f:
        f.write("이것은 샘플 텍스트 파일입니다.\n")
        f.write("테스트 목적으로 일부 샘플 텍스트가 포함되어 있습니다.")

def main():
    if not os.path.exists("samples"):
        os.makedirs("samples")

    create_sample_pdf("samples/sample_ko.pdf")
    create_sample_ppt("samples/sample_ko.pptx")
    create_sample_excel("samples/sample_ko.xlsx")
    create_sample_txt("samples/sample_ko.txt")
    
    pdf_text = extract_text_from_pdf("samples/sample_ko.pdf")
    ppt_text = extract_text_from_ppt("samples/sample_ko.pptx")
    
    insert_into_database("샘플 PDF", pdf_text)
    insert_into_database("샘플 PPT", ppt_text)
    insert_into_database("샘플 엑셀", "샘플 엑셀 파일\n테스트 목적으로 일부 샘플 텍스트가 포함되어 있습니다.")
    insert_into_database("샘플 텍스트", "이것은 샘플 텍스트 파일입니다.\n테스트 목적으로 일부 샘플 텍스트가 포함되어 있습니다.")

if __name__ == "__main__":
    main()
