import sys
import os

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..')))

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from backend.utils.insert_data import insert_into_database
from backend.utils.extract_text import extract_text_from_pdf, extract_text_from_ppt

def create_sample_pdf(file_path):
    c = canvas.Canvas(file_path, pagesize=letter)
    c.drawString(100, 750, "This is a sample PDF document.")
    c.drawString(100, 735, "It contains some sample text for testing purposes.")
    c.save()

def create_sample_ppt(file_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sample Presentation"
    subtitle.text = "This is a sample PPT slide for testing purposes."
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title, content = slide.shapes.title, slide.placeholders[1]
    title.text = "Slide 2 Title"
    content.text = "This slide contains some more sample text."
    prs.save(file_path)

def main():
    create_sample_pdf("sample.pdf")
    create_sample_ppt("sample.pptx")
    pdf_text = extract_text_from_pdf("sample.pdf")
    ppt_text = extract_text_from_ppt("sample.pptx")
    insert_into_database("Sample PDF", pdf_text)
    insert_into_database("Sample PPT", ppt_text)

if __name__ == "__main__":
    main()
