import os

def create_project_structure(base_path):
    dirs = [
        "project_root/backend",
        "project_root/backend/db",
        "project_root/backend/utils",
        "project_root/frontend",
        "project_root/frontend/src",
        "project_root/frontend/public",
        "project_root/frontend/src/components"
    ]

    for dir in dirs:
        os.makedirs(os.path.join(base_path, dir), exist_ok=True)

    backend_files = {
        "project_root/backend/main.py": """
import sys
import os
from fastapi import FastAPI, File, UploadFile, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
import logging
import nltk
from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
from fuzzywuzzy import fuzz
import shutil
import pandas as pd
from gensim.models import Word2Vec

sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from backend.db.connection import create_connection
from backend.utils.extract_text import extract_text_from_pdf, extract_text_from_ppt, extract_text_from_excel
from backend.utils.insert_data import insert_into_database

# NLTK 데이터 다운로드 (최초 실행 시 한 번 필요)
nltk.download('stopwords')
nltk.download('wordnet')

app = FastAPI()

# CORS 설정
origins = [
    "http://localhost:3000",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

lemmatizer = WordNetLemmatizer()
stop_words = set(stopwords.words('english'))

def preprocess_text(text):
    words = text.lower().split()
    words = [lemmatizer.lemmatize(word) for word in words if word not in stop_words]
    return ' '.join(words)

def find_similar_sentences(query, documents, threshold=80):
    preprocessed_query = preprocess_text(query)
    matches = []
    for doc in documents:
        sentences = doc['content'].split('. ')
        for sentence in sentences:
            preprocessed_sentence = preprocess_text(sentence)
            if fuzz.partial_ratio(preprocessed_query, preprocessed_sentence) >= threshold:
                matches.append({
                    'id': doc['id'],
                    'title': doc['title'],
                    'content': sentence
                })
    return matches

@app.get("/search")
async def search_documents(query: str):
    logger.info(f"Received search request: {query}")
    connection = create_connection()
    cursor = connection.cursor(dictionary=True)
    
    cursor.execute("SELECT id, title, LEFT(content, 200) AS content FROM documents")
    documents = cursor.fetchall()
    
    similar_sentences = find_similar_sentences(query, documents)
    
    cursor.close()
    connection.close()
    return similar_sentences

@app.post("/uploadfile/")
async def upload_file(file: UploadFile = File(...)):
    if not os.path.exists('files'):
        os.makedirs('files')

    file_location = f"files/{file.filename}"
    with open(file_location, "wb+") as file_object:
        shutil.copyfileobj(file.file, file_object)

    file_extension = file.filename.split('.')[-1]
    
    if file_extension == 'pdf':
        text = extract_text_from_pdf(file_location)
    elif file_extension in ['ppt', 'pptx']:
        text = extract_text_from_ppt(file_location)
    elif file_extension in ['xlsx', 'xls']:
        text = extract_text_from_excel(file_location)
    elif file_extension == 'txt':
        with open(file_location, 'r', encoding='utf-8') as f:
            text = f.read()
    else:
        return {"error": "Unsupported file type"}

    insert_into_database(file.filename, text)
    return {"info": "File uploaded successfully"}

@app.get("/document/{document_id}")
async def get_document(document_id: int):
    connection = create_connection()
    cursor = connection.cursor(dictionary=True)
    cursor.execute("SELECT * FROM documents WHERE id = %s", (document_id,))
    document = cursor.fetchone()
    cursor.close()
    connection.close()
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    return document

@app.delete("/document/{document_id}")
async def delete_document(document_id: int):
    connection = create_connection()
    cursor = connection.cursor()
    cursor.execute("SELECT * FROM documents WHERE id = %s", (document_id,))
    document = cursor.fetchone()
    if not document:
        cursor.close()
        connection.close()
        raise HTTPException(status_code=404, detail="Document not found")

    # Delete the file from the filesystem
    file_path = f"files/{document[1]}"  # Assuming the file name is in the second column
    if os.path.exists(file_path):
        os.remove(file_path)

    # Delete the record from the database
    cursor.execute("DELETE FROM documents WHERE id = %s", (document_id,))
    connection.commit()
    
    cursor.close()
    connection.close()
    return {"info": "Document deleted successfully"}
""",
        "project_root/backend/db/connection.py": """
import mysql.connector

def create_connection():
    connection = mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="documents_db"
    )
    return connection
""",
        "project_root/backend/utils/extract_text.py": """
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation

def extract_text_from_pdf(file_path):
    document = fitz.open(file_path)
    text = ""
    for page_num in range(document.page_count):
        page = document.load_page(page_num)
        text += page.get_text()
    return text

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()
""",
        "project_root/backend/utils/insert_data.py": """
import mysql.connector

def insert_into_database(file_name, content):
    connection = mysql.connector.connect(
        host="localhost",
        user="root",
        password="",
        database="documents_db"
    )
    cursor = connection.cursor()
    cursor.execute("INSERT INTO documents (title, content) VALUES (%s, %s)", (file_name, content))
    connection.commit()
    cursor.close()
    connection.close()
"""
    }

    for file_path, content in backend_files.items():
        with open(os.path.join(base_path, file_path), 'w', encoding='utf-8') as file:
            file.write(content.strip())

    frontend_files = {
        "project_root/frontend/src/App.js": """
import React, { useState } from 'react';
import './App.css';

function App() {
  const [query, setQuery] = useState('');
  const [useRelatedSearch, setUseRelatedSearch] = useState(false);
  const [results, setResults] = useState([]);
  const [noResults, setNoResults] = useState(false);
  const [loading, setLoading] = useState(false);
  const [error, setError] = useState('');
  const [file, setFile] = useState(null);
  const [selectedDocument, setSelectedDocument] = useState(null);

  const handleSearch = async () => {
    setLoading(true);
    setError('');
    setNoResults(false);
    setResults([]);
    try {
      const response = await fetch(`http://localhost:8000/search?query=${query}&related=${useRelatedSearch}`);
      if (!response.ok) {
        throw new Error('Failed to fetch results');
      }
      const data = await response.json();
      setResults(data);
      setNoResults(data.length === 0);
    } catch (err) {
      setError('An error occurred during the search. Please try again.');
    } finally {
      setLoading(false);
    }
  };

  const handleFileChange = (e) => {
    setFile(e.target.files[0]);
  };

  const handleFileUpload = async () => {
    if (!file) {
      setError('Please select a file to upload.');
      return;
    }

    const formData = new FormData();
    formData.append('file', file);

    try {
      const response = await fetch('http://localhost:8000/uploadfile/', {
        method: 'POST',
        body: formData,
      });

      if (!response.ok) {
        throw new Error('Failed to upload file');
      }

      const data = await response.json();
      if (data.error) {
        setError(data.error);
      } else {
        setError('');
        alert('File uploaded successfully');
      }
    } catch (err) {
      setError('An error occurred during the file upload. Please try again.');
    }
  };

  const fetchDocument = async (id) => {
    try {
      const response = await fetch(`http://localhost:8000/document/${id}`);
      if (!response.ok) {
        throw new Error('Failed to fetch document');
      }
      const data = await response.json();
      setSelectedDocument(data);
    } catch (err) {
      setError('An error occurred while fetching the document. Please try again.');
    }
  };

  const deleteDocument = async (id)```jsx
    try {
      const response = await fetch(`http://localhost:8000/document/${id}`, {
        method: 'DELETE',
      });
      if (!response.ok) {
        throw new Error('Failed to delete document');
      }
      alert('Document deleted successfully');
      handleSearch(); // Refresh the search results
    } catch (err) {
      setError('An error occurred while deleting the document. Please try again.');
    }
  };

  return (
    <div className="container">
      <h1>Search Documents</h1>
      <div className="search-bar">
        <input
          type="text"
          value={query}
          onChange={(e) => setQuery(e.target.value)}
          className="search-input"
        />
        <button onClick={handleSearch} disabled={loading}>
          {loading ? 'Searching...' : 'Search'}
        </button>
      </div>
      <label>
        <input
          type="checkbox"
          checked={useRelatedSearch}
          onChange={(e) => setUseRelatedSearch(e.target.checked)}
        />
        Use related search
      </label>
      {error && <p className="error">{error}</p>}
      {noResults && <p className="no-results">No search results found.</p>}
      <div className="results">
        <ul>
          {results.map((result) => (
            <li key={result.id}>
              <div onClick={() => fetchDocument(result.id)}>
                {result.title}: {result.content}...
              </div>
              <button onClick={() => deleteDocument(result.id)}>Delete</button>
            </li>
          ))}
        </ul>
      </div>
      {selectedDocument && (
        <div className="document-details">
          <h2>{selectedDocument.title}</h2>
          <p>{selectedDocument.content}</p>
          <button onClick={() => setSelectedDocument(null)}>Close</button>
        </div>
      )}
      <h2>Upload File</h2>
      <div className="upload-section">
        <label htmlFor="file-upload" className="custom-file-upload">
          {file ? file.name : 'Choose File'}
        </label>
        <input id="file-upload" type="file" onChange={handleFileChange} />
        <button onClick={handleFileUpload}>Upload</button>
      </div>
    </div>
  );
}

export default App;
""",
        "project_root/frontend/src/App.css": """
.container {
  text-align: center;
  max-width: 800px;
  margin: 0 auto;
  padding: 20px;
}

.search-bar {
  display: flex;
  justify-content: center;
  margin-bottom: 20px;
}

.search-input {
  width: 60%;
  padding: 10px;
  font-size: 16px;
}

button {
  padding: 10px 20px;
  font-size: 16px;
  margin-left: 10px;
}

.results {
  max-height: 300px;
  overflow-y: auto;
  margin-bottom: 20px;
}

.no-results, .error {
  color: red;
}

.document-details {
  border: 1px solid #ccc;
  padding: 10px;
  margin-top: 20px;
  max-height: 300px;
  overflow-y: auto;
}

.upload-section {
  display: flex;
  justify-content: center;
  margin-top: 20px;
}

.custom-file-upload {
  border: 1px solid #ccc;
  display: inline-block;
  padding: 10px 20px;
  cursor: pointer;
  margin-right: 10px;
}
""",
        "project_root/frontend/public/index.html": """
<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Document Search</title>
</head>
<body>
  <div id="root"></div>
</body>
</html>
""",
        "project_root/frontend/src/index.js": """
import React from 'react';
import ReactDOM from 'react-dom';
import './index.css';
import App from './App';

ReactDOM.render(
  <React.StrictMode>
    <App />
  </React.StrictMode>,
  document.getElementById('root')
);
"""
    }

    for file_path, content in frontend_files.items():
        with open(os.path.join(base_path, file_path), 'w', encoding='utf-8') as file:
            file.write(content.strip())

if __name__ == "__main__":
    base_path = "."  # 현재 디렉토리를 기준으로 생성
    create_project_structure(base_path)
    print("Project structure created successfully.")
